from transformers import T5Tokenizer, T5ForConditionalGeneration,pipeline
import wave,math,contextlib,nltk,sys,re,os,spacy,PyPDF2,bs4 as bs
from sklearn.feature_extraction.text import TfidfVectorizer
from sumy.summarizers.text_rank import TextRankSummarizer
from sumy.summarizers.lex_rank import LexRankSummarizer
from django.views.generic import TemplateView,FormView
from django.core.files.storage import default_storage
from sumy.parsers.plaintext import PlaintextParser
from .nltk_summarization import nltk_summarizer
from sklearn.decomposition import TruncatedSVD
from sklearn.pipeline import make_pipeline
from sumy.nlp.tokenizers import Tokenizer
from moviepy.editor import AudioFileClip
from nltk.stem import WordNetLemmatizer 
from django.shortcuts import render
from urllib.request import urlopen
from django.conf import settings
import speech_recognition as sr
from pydub import AudioSegment
from textsummy.models import *
from bs4 import BeautifulSoup
from pptx.util import Pt
import urllib.request
from .forms import *
import pptx
import docx



#Execute this line if you are running this code for first time
nltk.download('wordnet')

#Initializing few variable
nlp = spacy.load('en_core_web_sm')
lemmatizer = WordNetLemmatizer() 

# Sumy 
def sumy_doc_summary(docx):
	parser = PlaintextParser.from_string(docx,Tokenizer("english"))
	lex_summarizer = LexRankSummarizer()
	summary = lex_summarizer(parser.document,3)
	summary_list = [str(sentence) for sentence in summary]
	result = ' '.join(summary_list)
	return result

def summy_summary(text):
    parser = PlaintextParser.from_string(text, Tokenizer('english'))
    summarizer = TextRankSummarizer()
    summarizer.stop_words = ['in', 'the', 'and', 'but', 'a', 'an']
    summary = summarizer(parser.document, sentences_count=10)
    summary_text = [str(sentence) for sentence in summary]
    return ' '.join(summary_text)
 
def spaacy_summary(input_text, num_sentences=30):
    # Load the Spacy model
    nlp = spacy.load('en_core_web_sm')

    # Process the input text
    doc = nlp(input_text)

    # Get the sentences and their corresponding scores
    sentence_scores = {}
    for sent in doc.sents:
        sentence_scores[sent] = sent.similarity(doc)

    # Sort the sentences by score and get the top `num_sentences`
    top_sentences = sorted(sentence_scores, key=sentence_scores.get, reverse=True)[:num_sentences]

    # Join the top sentences to create the summary
    summary = ' '.join([sent.text for sent in top_sentences])

    return summary
 
# Reading Time
def readingTime(mytext):
	total_words = len([ token.text for token in nlp(mytext)])
	estimatedTime = total_words/200.0
	return estimatedTime

# Fetch Text From Url
def get_text(url):
	page = urlopen(url)
	soup = BeautifulSoup(page)
	fetched_text = ' '.join(map(lambda p:p.text,soup.find_all('p')))
	return fetched_text

# Define a function to generate a summary
def generate_summary(text):
    # Initialize the model and tokenizer
    model_name = 't5-small'
    tokenizer = T5Tokenizer.from_pretrained(model_name)
    model = T5ForConditionalGeneration.from_pretrained(model_name)
    input_ids = tokenizer.encode(text, return_tensors='pt')

    # Generate a summary
    summary_ids = model.generate(input_ids, num_beams=4, max_length=100, early_stopping=True)
    summary = tokenizer.decode(summary_ids[0], skip_special_tokens=True)
    return summary

def lsa(text, n=1):
    # Create a pipeline to extract the latent topics
    pipeline = make_pipeline(
        TfidfVectorizer(),
        TruncatedSVD(n_components=n)
    )

    # Fit the pipeline to the text data and extract the top n topics
    pipeline.fit_transform([text])
    topics = pipeline.named_steps['truncatedsvd'].components_

    # Create the summary by selecting the top sentences based on the most important topics
    sentences = text.split('.')
    sentence_scores = []
    for i in range(len(sentences)):
        sentence_score = 0
        for j in range(n):
            sentence_score += abs(topics[j][i])
        sentence_scores.append((i, sentence_score))
    top_sentences = sorted(sentence_scores, key=lambda item: item[1], reverse=True)[:n]

    # Create the
    summary = ''
    for i in sorted([x[0] for x in top_sentences]):
        summary += sentences[i] + '.'

    return summary

def get_text_from_pdf(pdf_path):
    try:
        with open(pdf_path,'rb') as f:
            pdf = PyPDF2.PdfReader(f)
            text = ''
            for page in range(len(pdf.pages)):
                text += pdf.pages[page].extract_text()
        return text
    finally:
        os.remove(pdf_path)

def get_text_from_file(file_path):
    try:
        with open(file_path) as f:
            text = f.read().replace("\n", '')
        return text
    finally:
        os.remove(file_path)

def get_text_from_site(url):
    scrap_data = urllib.request.urlopen(url)
    article = scrap_data.read()
    parsed_article = bs.BeautifulSoup(article,'lxml')
    
    paragraphs = parsed_article.find_all('p')
    article_text = ""
    
    for p in paragraphs:
        article_text += p.text
    
    #Removing all unwanted characters
    article_text = re.sub(r'\[[0-9]*\]', '', article_text)
    return article_text

def get_text_from_video(video_path):
    try:
        transcribed_audio_file_name = "media\\wavs\\transcribed_speech.wav"
        audioclip = AudioFileClip(video_path)
        audioclip.write_audiofile(transcribed_audio_file_name)
        with contextlib.closing(wave.open(transcribed_audio_file_name,'r')) as f:
            frames = f.getnframes()
            rate = f.getframerate()
            duration = frames / float(rate)
        total_duration = math.ceil(duration / 60)
        r = sr.Recognizer()
        text = ""
        for i in range(0, total_duration):
            with sr.AudioFile(transcribed_audio_file_name) as source:
                audio = r.record(source, offset=i*60, duration=60)
            text += r.recognize_google(audio)
        return text
    finally:
        os.remove(transcribed_audio_file_name)
        os.remove(video_path)

def get_text_from_docx(doc_path):
    try:
        document = docx.Document(doc_path)
        text = ""
        for paragraph in document.paragraphs:
            text += paragraph.text    
        return text
    finally:
       os.remove(doc_path)

def get_text_from_ppt(ppt_path):
    try:
        prs = pptx.Presentation(ppt_path)
        text = ''
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text_frame = shape.text_frame
                    for paragraph in text_frame.paragraphs:
                        for run in paragraph.runs:
                            text += run.text
        return text
    finally:
        os.remove(ppt_path)

def frequency_matrix(sentences):
    freq_matrix = {}
    stopWords = nlp.Defaults.stop_words

    for sent in sentences:
        freq_table = {} #dictionary with 'words' as key and their 'frequency' as value
        
        #Getting all word from the sentence in lower case
        words = [word.text.lower() for word in sent  if word.text.isalnum()]
       
        for word in words:  
            word = lemmatizer.lemmatize(word)   #Lemmatize the word
            if word not in stopWords:           #Reject stopwords
                if word in freq_table:
                    freq_table[word] += 1
                else:
                    freq_table[word] = 1

        freq_matrix[sent[:15]] = freq_table

    return freq_matrix

def tf_matrix(freq_matrix):
    tf_matrix = {}

    for sent, freq_table in freq_matrix.items():
        tf_table = {}  #dictionary with 'word' itself as a key and its TF as value

        total_words_in_sentence = len(freq_table)
        for word, count in freq_table.items():
            tf_table[word] = count / total_words_in_sentence

        tf_matrix[sent] = tf_table

    return tf_matrix

def sentences_per_words(freq_matrix):
    sent_per_words = {}

    for sent, f_table in freq_matrix.items():
        for word, count in f_table.items():
            if word in sent_per_words:
                sent_per_words[word] += 1
            else:
                sent_per_words[word] = 1

    return sent_per_words

def idf_matrix(freq_matrix, sent_per_words, total_sentences):
    idf_matrix = {}

    for sent, f_table in freq_matrix.items():
        idf_table = {}

        for word in f_table.keys():
            idf_table[word] = math.log10(total_sentences / float(sent_per_words[word]))

        idf_matrix[sent] = idf_table

    return idf_matrix

def tf_idf_matrix(tf_matrix, idf_matrix):
    tf_idf_matrix = {}

    for (sent1, f_table1), (sent2, f_table2) in zip(tf_matrix.items(), idf_matrix.items()):

        tf_idf_table = {}

       #word1 and word2 are same
        for (word1, tf_value), (word2, idf_value) in zip(f_table1.items(),
                                                    f_table2.items()):  
            tf_idf_table[word1] = float(tf_value * idf_value)

        tf_idf_matrix[sent1] = tf_idf_table

    return tf_idf_matrix

def score_sentences(tf_idf_matrix):
    
    sentenceScore = {}

    for sent, f_table in tf_idf_matrix.items():
        total_tfidf_score_per_sentence = 0

        total_words_in_sentence = len(f_table)
        for word, tf_idf_score in f_table.items():
            total_tfidf_score_per_sentence += tf_idf_score

        if total_words_in_sentence != 0:
            sentenceScore[sent] = total_tfidf_score_per_sentence / total_words_in_sentence

    return sentenceScore

def average_score(sentence_score):
    
    total_score = 0
    for sent in sentence_score:
        total_score += sentence_score[sent]

    average_sent_score = (total_score / len(sentence_score))

    return average_sent_score

def create_summary(sentences, sentence_score, threshold):
    summary = ''

    for sentence in sentences:
        if sentence[:15] in sentence_score and sentence_score[sentence[:15]] >= (threshold):
            summary += " " + sentence.text
        

    return summary

# Splash Screen Page <BACK-END/>
class HomePageView(TemplateView):
    template_name = 'home.html'

# Typing text or recording voice Summary Page <BACK-END/>
class TextSummaryPageView(FormView):
    template_name = 'text_summary.html'
    success_url = '/text-summary/'
    form_class = TextSummaryForm

    def form_valid(self, form):
        text = form.cleaned_data['text_summary']
        txt = nlp(text)
        sentences = list(txt.sents)
        total_sentences = len(sentences)
        freq_matrix = frequency_matrix(sentences)
        tf_matrixes = tf_matrix(freq_matrix)
        num_sent_per_words = sentences_per_words(freq_matrix)
        idf_matrixes = idf_matrix(freq_matrix, num_sent_per_words, total_sentences)
        tf_idf_matrixes = tf_idf_matrix(tf_matrixes, idf_matrixes)
        sentence_scores = score_sentences(tf_idf_matrixes)
        threshold = average_score(sentence_scores)
        summary = create_summary(sentences, sentence_scores, 1.3 * threshold)
        nltk_summary = nltk_summarizer(text)
        sumy_summary = summy_summary(text)
        spacy_summary = spaacy_summary(text)
        t5_summary   = generate_summary(text)
        lsa_summary  = lsa(text)
        context      = self.get_context_data(form=form, 
                                             spacy_summary=spacy_summary,
                                             nltk_summary=summary if len(summary) > 10 else nltk_summary,
                                             t5_summary=t5_summary,
                                             sumy_summary=sumy_summary,
                                             lsa_summary=lsa_summary
                                             )
        return self.render_to_response(context)

# Files Summary Page <BACK-END/>
class FilesSummaryPageView(FormView):
    template_name = 'files_summary.html'
    success_url = '/files-summary/'
    form_class = FilesSummaryForm

    def form_valid(self, form):
        file = form.cleaned_data['file_summary']
        file_path = default_storage.save(file.name, file)
        file_url = self.request.scheme + '://' + self.request.get_host() + settings.MEDIA_URL + file_path
        if '.pdf' in file_path:
            text = get_text_from_pdf(f"media\\{file_path}")
        elif '.txt' in file_path:
            text = get_text_from_file(f"media\\{file_path}")
        elif '.docx' in file_path:
            text = get_text_from_docx(f"media\\{file_path}")
        elif '.pptx' in file_path:
            text = get_text_from_ppt(f"media\\{file_path}")
        else:
            pass
        txt = nlp(text)
        sentences = list(txt.sents)
        total_sentences = len(sentences)
        freq_matrix = frequency_matrix(sentences)
        tf_matrixes = tf_matrix(freq_matrix)
        num_sent_per_words = sentences_per_words(freq_matrix)
        idf_matrixes = idf_matrix(freq_matrix, num_sent_per_words, total_sentences)
        tf_idf_matrixes = tf_idf_matrix(tf_matrixes, idf_matrixes)
        sentence_scores = score_sentences(tf_idf_matrixes)
        threshold = average_score(sentence_scores)
        summary = create_summary(sentences, sentence_scores, 1.3 * threshold)
        nltk_summary = nltk_summarizer(text)
        sumy_summary = summy_summary(text)
        spacy_summary = spaacy_summary(text)
        t5_summary   = generate_summary(text)
        lsa_summary  = lsa(text)
        context      = self.get_context_data(form=form, 
                                             spacy_summary=spacy_summary,
                                             nltk_summary=summary if len(summary) > 10 else nltk_summary,
                                             t5_summary=t5_summary,
                                             sumy_summary=sumy_summary,
                                             lsa_summary=lsa_summary
                                             )
        return self.render_to_response(context)

# Video Summary Page <BACK-END/>
class VideoSummaryPageView(FormView):
    template_name = 'video_summary.html'
    success_url = '/video-summary/'
    form_class = VideoSummaryForm
    
    def form_valid(self, form):
        video = form.cleaned_data['video_summary']
        video_path = default_storage.save(video.name, video)
        video_url = self.request.scheme + '://' + self.request.get_host() + settings.MEDIA_URL + video_path
        text = get_text_from_video(f"media\\{video_path}")
        txt = nlp(text)
        sentences = list(txt.sents)
        total_sentences = len(sentences)
        freq_matrix = frequency_matrix(sentences)
        tf_matrixes = tf_matrix(freq_matrix)
        num_sent_per_words = sentences_per_words(freq_matrix)
        idf_matrixes = idf_matrix(freq_matrix, num_sent_per_words, total_sentences)
        tf_idf_matrixes = tf_idf_matrix(tf_matrixes, idf_matrixes)
        sentence_scores = score_sentences(tf_idf_matrixes)
        threshold = average_score(sentence_scores)
        summary = create_summary(sentences, sentence_scores, 1.3 * threshold)
        nltk_summary = nltk_summarizer(text)
        sumy_summary = summy_summary(text)
        spacy_summary = spaacy_summary(text)
        t5_summary   = generate_summary(text)
        lsa_summary  = lsa(text)
        context      = self.get_context_data(form=form, 
                                             spacy_summary=spacy_summary,
                                             nltk_summary=summary if len(summary) > 10 else nltk_summary,
                                             t5_summary=t5_summary,
                                             sumy_summary=sumy_summary,
                                             lsa_summary=lsa_summary
                                             )
        return self.render_to_response(context)

# Website Summary Page <BACK-END/>
class UrlSummaryPageView(FormView):
    template_name = 'url_summary.html'
    success_url = '/url-summary/'
    form_class = UrlSummaryForm

    def form_valid(self, form):
        url           = form.cleaned_data.get('url_summary')
        text          = get_text_from_site(url) 
        txt = nlp(text)
        sentences = list(txt.sents)
        total_sentences = len(sentences)
        freq_matrix = frequency_matrix(sentences)
        tf_matrixes = tf_matrix(freq_matrix)
        num_sent_per_words = sentences_per_words(freq_matrix)
        idf_matrixes = idf_matrix(freq_matrix, num_sent_per_words, total_sentences)
        tf_idf_matrixes = tf_idf_matrix(tf_matrixes, idf_matrixes)
        sentence_scores = score_sentences(tf_idf_matrixes)
        threshold = average_score(sentence_scores)
        summary = create_summary(sentences, sentence_scores, 1.3 * threshold)
        nltk_summary  = nltk_summarizer(text)   
        sumy_summary  = summy_summary(text)
        spacy_summary = spaacy_summary(text)
        t5_summary    = generate_summary(text)
        lsa_summary   = lsa(text)
        context       = self.get_context_data(form=form, 
                                             spacy_summary=spacy_summary,
                                             nltk_summary=summary if len(summary) > 10 else nltk_summary,
                                             t5_summary=t5_summary,
                                             sumy_summary=sumy_summary,
                                             lsa_summary=lsa_summary)
        return self.render_to_response(context)
